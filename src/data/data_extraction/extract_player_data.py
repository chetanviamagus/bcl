#!/usr/bin/env python3
"""
Script to extract player details from PowerPoint file and create structured data files.
"""

import os
import re
import json
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

def clean_text(text):
    """Clean and normalize text content."""
    if not text:
        return ""
    # Remove extra whitespace and normalize line breaks
    text = re.sub(r'\s+', ' ', text.strip())
    return text

def extract_player_name(text):
    """Extract player name from text."""
    # Look for patterns like "Name: John Doe" or just "John Doe"
    if 'Name:' in text or 'NAME:' in text:
        name = text.split(':', 1)[1].strip()
    else:
        # If no "Name:" prefix, assume the first part is the name
        name = text.split('\n')[0].strip()
        name = name.split('Category:')[0].strip()
    
    return clean_text(name)

def extract_age(text):
    """Extract age from text."""
    age_match = re.search(r'(\d+)', text)
    return age_match.group(1) if age_match else ""

def extract_category(text):
    """Extract category from text."""
    if 'Category:' in text or 'CATEGORY:' in text:
        category = text.split(':', 1)[1].strip()
        # Clean up category text
        category = re.sub(r'[^\w\s-]', '', category).strip()
        return category
    return ""

def extract_mobile(text):
    """Extract mobile number from text."""
    mobile_match = re.search(r'(\d{10,})', text)
    return mobile_match.group(1) if mobile_match else ""

def extract_player_details(pptx_path):
    """Extract all player details from PowerPoint file."""
    
    if not os.path.exists(pptx_path):
        print(f"Error: PowerPoint file not found at {pptx_path}")
        return []
    
    try:
        # Load the presentation
        prs = Presentation(pptx_path)
        print(f"Loaded presentation with {len(prs.slides)} slides")
        
        players = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            # Look for text boxes that might contain player information
            text_content = []
            
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text.strip():
                    text_content.append(shape.text.strip())
            
            # Combine all text content
            full_text = ' '.join(text_content)
            
            # Try to extract player information
            if full_text and any(keyword in full_text.lower() for keyword in ['name', 'age', 'category', 'batsman', 'bowler', 'all rounder']):
                player_info = {
                    'slide_number': slide_num,
                    'name': extract_player_name(full_text),
                    'age': extract_age(full_text),
                    'category': extract_category(full_text),
                    'mobile': extract_mobile(full_text),
                    'raw_text': full_text
                }
                
                # Only add if we have at least name and age
                if player_info['name'] and player_info['age']:
                    players.append(player_info)
                    print(f"Slide {slide_num}: {player_info['name']} - {player_info['age']} - {player_info['category']}")
        
        return players
        
    except Exception as e:
        print(f"Error processing PowerPoint file: {e}")
        return []

def create_structured_files(players):
    """Create various structured data files."""
    
    # Create output directory
    output_dir = "/Users/chetan/Documents/CodeProjects/ReactProjects/bcl/src/data"
    os.makedirs(output_dir, exist_ok=True)
    
    # 1. JSON file for easy parsing
    json_file = os.path.join(output_dir, "players_data.json")
    with open(json_file, 'w', encoding='utf-8') as f:
        json.dump(players, f, indent=2, ensure_ascii=False)
    print(f"Created JSON file: {json_file}")
    
    # 2. CSV file for spreadsheet viewing
    csv_file = os.path.join(output_dir, "players_data.csv")
    with open(csv_file, 'w', encoding='utf-8') as f:
        f.write("Slide,Name,Age,Category,Mobile,Raw_Text\n")
        for player in players:
            f.write(f"{player['slide_number']},{player['name']},{player['age']},{player['category']},{player['mobile']},\"{player['raw_text']}\"\n")
    print(f"Created CSV file: {csv_file}")
    
    # 3. TypeScript interface file for React
    ts_file = os.path.join(output_dir, "players_data.ts")
    with open(ts_file, 'w', encoding='utf-8') as f:
        f.write("export interface PlayerData {\n")
        f.write("  slide_number: number;\n")
        f.write("  name: string;\n")
        f.write("  age: string;\n")
        f.write("  category: string;\n")
        f.write("  mobile: string;\n")
        f.write("  raw_text: string;\n")
        f.write("}\n\n")
        f.write("export const playersData: PlayerData[] = ")
        f.write(json.dumps(players, indent=2, ensure_ascii=False))
        f.write(";\n")
    print(f"Created TypeScript file: {ts_file}")
    
    # 4. Summary statistics
    stats_file = os.path.join(output_dir, "players_summary.txt")
    with open(stats_file, 'w', encoding='utf-8') as f:
        f.write("BCL 2024 Players Data Summary\n")
        f.write("=" * 40 + "\n\n")
        f.write(f"Total Players: {len(players)}\n\n")
        
        # Category breakdown
        categories = {}
        for player in players:
            cat = player['category'] or 'Unknown'
            categories[cat] = categories.get(cat, 0) + 1
        
        f.write("Category Breakdown:\n")
        for cat, count in sorted(categories.items()):
            f.write(f"  {cat}: {count} players\n")
        
        f.write(f"\nAge Range:\n")
        ages = [int(p['age']) for p in players if p['age'].isdigit()]
        if ages:
            f.write(f"  Youngest: {min(ages)} years\n")
            f.write(f"  Oldest: {max(ages)} years\n")
            f.write(f"  Average: {sum(ages)/len(ages):.1f} years\n")
        
        f.write(f"\nPlayers with Mobile Numbers: {sum(1 for p in players if p['mobile'])}\n")
        f.write(f"Players without Mobile Numbers: {sum(1 for p in players if not p['mobile'])}\n")
    
    print(f"Created summary file: {stats_file}")
    
    return {
        'json': json_file,
        'csv': csv_file,
        'typescript': ts_file,
        'summary': stats_file
    }

def main():
    # Paths
    pptx_path = "/Users/chetan/Documents/CodeProjects/ReactProjects/bcl/src/assets/New-2024-BCL-Players.key"
    
    print("BCL Player Data Extractor")
    print("=" * 40)
    print(f"PowerPoint file: {pptx_path}")
    print()
    
    # Extract player details
    players = extract_player_details(pptx_path)
    
    if players:
        print(f"\nExtracted {len(players)} players")
        
        # Create structured files
        files = create_structured_files(players)
        
        print(f"\nFiles created:")
        for file_type, file_path in files.items():
            print(f"  {file_type.upper()}: {file_path}")
    else:
        print("No players found!")

if __name__ == "__main__":
    main()
