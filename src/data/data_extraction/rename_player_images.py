#!/usr/bin/env python3
"""
Script to rename player images with proper naming convention.
This script will rename the extracted images to follow the pattern: Name_Age_Category_MobileNumber.jpg
"""

import os
import re
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

def clean_filename(name, age, category, mobile=""):
    """Clean and format the filename according to the specified pattern."""
    # Remove special characters and spaces, replace with underscores
    clean_name = re.sub(r'[^\w\s-]', '', name).strip()
    clean_name = re.sub(r'[-\s]+', '_', clean_name)
    
    clean_category = re.sub(r'[^\w\s-]', '', category).strip()
    clean_category = re.sub(r'[-\s]+', '_', clean_category)
    
    # Clean mobile number (remove spaces, dashes, etc.)
    if mobile:
        clean_mobile = re.sub(r'[^\d]', '', mobile)
    else:
        clean_mobile = "0000000000"  # Default if no mobile found
    
    # Format: Name_Age_Category_MobileNumber
    filename = f"{clean_name}_{age}_{clean_category}_{clean_mobile}.jpg"
    return filename

def extract_and_rename_images(pptx_path, players_dir):
    """Extract player information and rename images accordingly."""
    
    if not os.path.exists(pptx_path):
        print(f"Error: PowerPoint file not found at {pptx_path}")
        return
    
    if not os.path.exists(players_dir):
        print(f"Error: Players directory not found at {players_dir}")
        return
    
    try:
        # Load the presentation
        prs = Presentation(pptx_path)
        print(f"Loaded presentation with {len(prs.slides)} slides")
        
        renamed_count = 0
        
        for slide_num, slide in enumerate(prs.slides, 1):
            # Look for text boxes that might contain player information
            player_info = {}
            text_content = []
            
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text.strip():
                    text_content.append(shape.text.strip())
            
            # Try to extract player information from text
            for text in text_content:
                # Look for patterns like "Name: John Doe", "Age: 25", etc.
                if 'Name:' in text or 'NAME:' in text:
                    player_info['name'] = text.split(':', 1)[1].strip()
                elif 'Age:' in text or 'AGE:' in text:
                    age_match = re.search(r'(\d+)', text)
                    if age_match:
                        player_info['age'] = age_match.group(1)
                elif 'Category:' in text or 'CATEGORY:' in text:
                    player_info['category'] = text.split(':', 1)[1].strip()
                elif 'Phone:' in text or 'PHONE:' in text or 'Mobile:' in text or 'MOBILE:' in text:
                    mobile_match = re.search(r'(\d{10,})', text)
                    if mobile_match:
                        player_info['mobile'] = mobile_match.group(1)
            
            # Check if we have enough information to rename
            if all(key in player_info for key in ['name', 'age', 'category']):
                # Find the corresponding image file
                old_filename = f"player_slide_{slide_num}_{slide_num}.jpg"
                old_path = os.path.join(players_dir, old_filename)
                
                if os.path.exists(old_path):
                    # Generate new filename
                    new_filename = clean_filename(
                        player_info['name'],
                        player_info['age'],
                        player_info['category'],
                        player_info.get('mobile', '')
                    )
                    new_path = os.path.join(players_dir, new_filename)
                    
                    # Rename the file
                    try:
                        os.rename(old_path, new_path)
                        print(f"Renamed: {old_filename} -> {new_filename}")
                        renamed_count += 1
                    except Exception as e:
                        print(f"Error renaming {old_filename}: {e}")
                else:
                    print(f"Image file not found: {old_filename}")
            else:
                print(f"Slide {slide_num}: Insufficient player info - {player_info}")
        
        print(f"\nRenaming complete! {renamed_count} images renamed successfully.")
        
    except Exception as e:
        print(f"Error processing PowerPoint file: {e}")

def main():
    # Paths
    pptx_path = "/Users/chetan/Documents/CodeProjects/ReactProjects/bcl/src/assets/New-2024-BCL-Players.key"
    players_dir = "/Users/chetan/Documents/CodeProjects/ReactProjects/bcl/src/assets/players"
    
    print("BCL Player Image Renamer")
    print("=" * 40)
    print(f"PowerPoint file: {pptx_path}")
    print(f"Players directory: {players_dir}")
    print()
    
    extract_and_rename_images(pptx_path, players_dir)

if __name__ == "__main__":
    main()
